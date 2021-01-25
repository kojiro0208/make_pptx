# pip install Pillow
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from PIL import Image

IMG_PATH = "../data/in/スマートグラス.jpg"
IMG_PATH2 = "../data/in/ヨガ２.jpg"


TITLE_SLIDE = prs.slide_layouts[0]
FREE_SLIDE = prs.slide_layouts[1]
BRANK_SLIDE = prs.slide_layouts[6]

#スライドオブジェクトの定義
prs = Presentation() 

slide = prs.slides.add_slide(TITLE_SLIDE)
add_title(slide, 'My slide' )

slide = prs.slides.add_slide(FREE_SLIDE)
add_title(slide, 'これは２枚目' )

#スライドサイズを取得

IMG_PATH = "../data/in/スマートグラス.jpg"
add_image(slide, IMG_PATH, right = 0.5)
IMG_PATH2 = "../data/in/ヨガ２.jpg"
add_image(slide, IMG_PATH2,  left = 0.5)

#白紙のスライドを追加
slide = prs.slides.add_slide(BRANK_SLIDE)
add_image(slide, IMG_PATH, right = 0.6, bottom = 0.3)
add_image(slide, IMG_PATH2, right = 0.5, top = 0.6)
add_image(slide, IMG_PATH, left = 0.5,  bottom = 0.3)
add_image(slide, IMG_PATH2, left = 0.6, top = 0.6)

#白紙のスライドを追加
prs = Presentation() 
slide = prs.slides.add_slide(BRANK_SLIDE)


slide = prs.slides.add_slide(BRANK_SLIDE)

msg = "データサイエンスとは、ホゲホゲ、データソリューション事業とは、hogehogefadfafawfafadwffggaedgfagfafgafafasfasgfagfasgfasgasgasgasgasdgasgawg"
add_text(slide, msg, left=0.5, top = 0.6, font_size=30)
add_text(slide, msg, left=0.5, font_size=30)

prs.save('../data/out/test.pptx')


